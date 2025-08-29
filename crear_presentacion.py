#!/usr/bin/env python3
"""
Script para crear presentación PowerPoint sobre la utilidad del álgebra en Ingeniería de Software
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def crear_presentacion():
    # Crear presentación
    prs = Presentation()
    
    # Colores corporativos
    azul = RGBColor(46, 134, 193)  # #2E86C1
    gris = RGBColor(86, 101, 115)  # #566573
    
    # DIAPOSITIVA 1: PORTADA
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Layout de título
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "LA UTILIDAD DEL ÁLGEBRA Y LAS ECUACIONES EN INGENIERÍA DE SOFTWARE"
    title1.text_frame.paragraphs[0].font.size = Pt(32)
    title1.text_frame.paragraphs[0].font.color.rgb = azul
    
    subtitle_text = """Presentado por: [Tu nombre]
Carrera: Ingeniería de Software
Fecha: Diciembre 2024"""
    subtitle1.text = subtitle_text
    subtitle1.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle1.text_frame.paragraphs[0].font.color.rgb = gris
    
    # DIAPOSITIVA 2: VIDA COTIDIANA
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Layout de título y contenido
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "UTILIDAD DEL ÁLGEBRA EN LA VIDA COTIDIANA"
    title2.text_frame.paragraphs[0].font.color.rgb = azul
    
    contenido2 = """¿Por qué es útil el álgebra en mi vida cotidiana?

Como estudiante de ingeniería de software y programador junior:
• En el trabajo: Optimización de algoritmos y complejidad temporal
• En los estudios: Resolución de problemas de programación
• En proyectos: Desarrollo de aplicaciones con cálculos matemáticos

EJEMPLO: Optimización de memoria en aplicación web

Problema: ¿Cuántos usuarios máximos puede soportar la aplicación?
• Memoria base: 200 MB
• Memoria por usuario: 15 MB
• Memoria máxima: 2048 MB

Ecuación: 2048 = 200 + 15x
Solución: x = (2048 - 200) ÷ 15 = 123 usuarios máximos"""
    
    content2.text = contenido2
    content2.text_frame.paragraphs[0].font.size = Pt(16)
    
    # DIAPOSITIVA 3: ESTUDIOS UNIVERSITARIOS
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "UTILIDAD PARA CULMINAR ESTUDIOS UNIVERSITARIOS"
    title3.text_frame.paragraphs[0].font.color.rgb = azul
    
    contenido3 = """¿Por qué es esencial para mis estudios?

El álgebra es base para:
• Análisis de algoritmos: Complejidad temporal y espacial
• Estructuras de datos: Optimización de operaciones
• Bases de datos: Normalización y optimización
• IA: Machine learning y redes neuronales

EJEMPLO: Análisis de Merge Sort

Ecuación de recurrencia: T(n) = 2T(n/2) + n
Aplicando teorema maestro:
• a = 2, b = 2, f(n) = n
• log₂(2) = 1, f(n) = n¹
• Resultado: T(n) = Θ(n log n)

El algoritmo es eficiente para grandes conjuntos de datos."""
    
    content3.text = contenido3
    content3.text_frame.paragraphs[0].font.size = Pt(16)
    
    # DIAPOSITIVA 4: DESARROLLO PROFESIONAL
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "UTILIDAD EN EL DESARROLLO PROFESIONAL"
    title4.text_frame.paragraphs[0].font.color.rgb = azul
    
    contenido4 = """¿Cómo me ayudará como Full Stack Developer?

Aplicaciones profesionales:
• Backend: Optimización de consultas SQL y algoritmos
• Frontend: Cálculos responsive design y animaciones
• DevOps: Escalamiento automático y balanceadores
• ML: Implementación de modelos predictivos

EJEMPLO: Sistema de escalamiento automático

Para 3500 requests/min:
• Servidores mínimos = ⌈3500/1000⌉ = 4 servidores
• Con factor seguridad 20%: 4 × 1.2 = 5 servidores
• Costo mensual = 5 × $50 = $250/mes

Optimización: Minimizar C = n × 50, sujeto a restricciones de rendimiento"""
    
    content4.text = contenido4
    content4.text_frame.paragraphs[0].font.size = Pt(16)
    
    # DIAPOSITIVA 5: REFERENCIAS
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "REFERENCIAS"
    title5.text_frame.paragraphs[0].font.color.rgb = azul
    
    referencias = """1. Cormen, T. H., et al. (2022). Introduction to Algorithms (4th ed.). MIT Press.

2. Sedgewick, R., & Wayne, K. (2011). Algorithms (4th ed.). Addison-Wesley.

3. Skiena, S. S. (2020). The Algorithm Design Manual (3rd ed.). Springer.

4. Knuth, D. E. (1997). The Art of Computer Programming, Vol. 1. Addison-Wesley.

5. AWS Documentation (2024). Auto Scaling User Guide. Amazon Web Services.

6. Google Cloud (2024). Compute Engine Autoscaler. Google Cloud Platform.

7. Bentley, J. (2000). Programming Pearls (2nd ed.). Addison-Wesley.

8. McConnell, S. (2004). Code Complete (2nd ed.). Microsoft Press."""
    
    content5.text = referencias
    content5.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Guardar presentación
    prs.save('/workspace/Algebra_Ingenieria_Software.pptx')
    print("✅ Presentación creada exitosamente: Algebra_Ingenieria_Software.pptx")

if __name__ == "__main__":
    crear_presentacion()